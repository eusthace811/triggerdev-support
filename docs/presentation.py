#!/usr/bin/env python3
"""Generate the Trigger Support project presentation with Trigger.dev branding."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Trigger.dev brand colors (extracted from logo SVG gradients) ───
# "trigger" text + logo mark: #41FF54 → #E7FF52 (green to yellow-green)
# ".dev" text: #2563EB → #A855F7 (blue to purple)
BG_BLACK = RGBColor(0x09, 0x09, 0x0B)       # Trigger.dev site background
BG_CARD = RGBColor(0x13, 0x13, 0x17)        # Slightly lighter card surface
GREEN_1 = RGBColor(0x41, 0xFF, 0x54)        # Primary green (bottom of gradient)
GREEN_2 = RGBColor(0xE7, 0xFF, 0x52)        # Yellow-green (top of gradient)
GREEN_MID = RGBColor(0x94, 0xFF, 0x53)      # Midpoint of green gradient
BLUE = RGBColor(0x25, 0x63, 0xEB)           # ".dev" blue (bottom)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)         # ".dev" purple (top)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)          # Slightly warm white
GRAY_1 = RGBColor(0xA1, 0xA1, 0xAA)         # Muted text
GRAY_2 = RGBColor(0x71, 0x71, 0x7A)         # Even more muted
GRAY_3 = RGBColor(0x3F, 0x3F, 0x46)         # Divider lines / subtle shapes
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)       # Error / warning accent
AMBER = RGBColor(0xF5, 0x9E, 0x0B)          # Caution accent

# Font
FONT = "Helvetica Neue"
FONT_MONO = "Menlo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

TOTAL = 24


# ─── Helpers ───────────────────────────────────────────────────────

def set_bg(slide, color=BG_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def text(slide, l, t, w, h, txt, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, font=FONT, spacing=None):
    """Add a single text box. Returns the shape."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    if spacing is not None:
        p.space_after = spacing
    return box


def bullets(slide, items, l=Inches(1.2), t=Inches(2.3), w=Inches(11),
            size=19, color=WHITE, line_spacing=Pt(10), font=FONT):
    """Add a multi-line bullet text box."""
    box = slide.shapes.add_textbox(l, t, w, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Detect sub-items (indented with spaces)
        is_sub = item.startswith("   ")
        is_empty = item.strip() == ""

        if is_sub:
            p.text = item.strip()
            p.font.size = Pt(size - 2)
            p.font.color.rgb = GRAY_1
        elif is_empty:
            p.text = ""
            p.font.size = Pt(size - 4)
        else:
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = color

        p.font.name = font
        p.space_after = line_spacing
    return box


def slide_num(slide, n):
    text(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
         f"{n}/{TOTAL}", size=10, color=GRAY_2, align=PP_ALIGN.RIGHT)


def logo(slide, l=Inches(0.5), t=Inches(0.35), h=Inches(0.42)):
    """Place the Trigger.dev logo PNG."""
    # Logo aspect ratio is 544:96 (5.667:1). Ensure it stays within slide bounds.
    logo_w = h * (544 / 96)  # auto width from height
    max_right = prs.slide_width - Inches(0.3)  # right margin
    if l + logo_w > max_right:
        l = max_right - logo_w
    slide.shapes.add_picture("docs/logo.png", l, t, height=h)


def section_header(slide, title, subtitle=None, n=None):
    """Standard section slide with title, green accent bar, optional subtitle."""
    set_bg(slide)
    # Thin green accent bar
    rect(slide, Inches(1.2), Inches(1.85), Inches(1.8), Pt(3), GREEN_1)
    # Title
    text(slide, Inches(1.2), Inches(0.7), Inches(10.5), Inches(1.0),
         title, size=38, color=WHITE, bold=True)
    if subtitle:
        text(slide, Inches(1.2), Inches(2.05), Inches(10.5), Inches(0.5),
             subtitle, size=18, color=GRAY_1)
    logo(slide, Inches(11.0), Inches(0.35), Inches(0.35))
    if n:
        slide_num(slide, n)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ═══════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════

# ── 1. TITLE ──────────────────────────────────────────────────────
s = blank()
set_bg(s)

# Top accent line (green gradient feel: green bar then yellow-green bar)
rect(s, 0, 0, Inches(8), Pt(4), GREEN_1)
rect(s, Inches(8), 0, Inches(5.333), Pt(4), GREEN_2)

# Logo large
s.shapes.add_picture("docs/logo.png", Inches(1.2), Inches(1.6), height=Inches(0.7))

# Title
text(s, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
     "AI Support Agent", size=52, color=WHITE, bold=True)
text(s, Inches(1.2), Inches(4.0), Inches(10), Inches(0.7),
     "Accurate, cited answers from 300+ documentation pages", size=22, color=GRAY_1)

# Subtitle info
text(s, Inches(1.2), Inches(5.3), Inches(10), Inches(0.5),
     "Juan Suarez  \u2022  April 2026  \u2022  Job Application Project", size=16, color=GRAY_2)

# Bottom accent
rect(s, 0, H - Pt(4), Inches(5), Pt(4), BLUE)
rect(s, Inches(5), H - Pt(4), Inches(8.333), Pt(4), PURPLE)

# Mascot
s.shapes.add_picture("docs/triggerito.png", Inches(10), Inches(2.2), height=Inches(3.0))

slide_num(s, 1)


# ── 2. AGENDA ─────────────────────────────────────────────────────
s = blank()
section_header(s, "Agenda", n=2)
bullets(s, [
    "1.  The Challenge",
    "2.  Dry Run: Getting Started with chat.agent()",
    "3.  Architecture & Tech Stack",
    "4.  Implementation Phases",
    "5.  Prerelease SDK Bugs & Workarounds",
    "6.  RAG Quality: Evaluation & Improvement Cycles",
    "7.  Results: 5.8 \u2192 8.8 \u2192 9.3 / 10",
    "8.  Production Deployment",
    "9.  Honest Gaps: What's NOT Production-Ready",
    "10. Roadmap to World-Class Support Agent",
    "11. Developer Experience Feedback",
], size=21, t=Inches(2.5))


# ── 3. THE PROBLEM ────────────────────────────────────────────────
s = blank()
section_header(s, "The Challenge", "Build a support agent that actually works", 3)

bullets(s, [
    "The brief:",
    '   "Build a Trigger.dev support agent using chat.agent() that answers',
    '    questions accurately. Walk us through your plan, what you built,',
    '    and how you\'d turn it into a world-class agent."',
    "",
    "What makes this hard:",
    "   300+ MDX docs pages covering tasks, scheduling, concurrency, realtime, ...",
    "   Developers ask natural language that doesn't map to SDK terminology",
    "   Prerelease SDK with incomplete docs and undiscovered bugs",
    "   5-hour window \u2014 need to scope aggressively",
    "",
    "My interpretation:",
    "   Ship something that genuinely works, not a toy demo",
    "   Show the thinking, not just the code",
    "   Be honest about what's missing and what comes next",
], t=Inches(2.5))


# ── 4. INITIAL THOUGHTS ──────────────────────────────────────────
s = blank()
section_header(s, "Initial Thoughts", "Two approaches considered", 4)

# Left column: Option A
text(s, Inches(1.2), Inches(2.5), Inches(5), Inches(0.5),
     "Option A: On-demand fetch", size=22, color=RED_SOFT, bold=True)
bullets(s, [
    "MCP servers / inline functions (fetchDoc, listDocs)",
    "Simple but slow, no semantic search, poor recall",
    "Every question = multiple API calls to GitHub",
], l=Inches(1.2), t=Inches(3.2), w=Inches(5), size=17)

# Right column: Option B
text(s, Inches(7.2), Inches(2.5), Inches(5), Inches(0.5),
     "Option B: Full RAG pipeline", size=22, color=GREEN_1, bold=True)
bullets(s, [
    "Pre-index all docs into a vector database",
    "Hybrid search (vector + BM25) for semantic + exact",
    "Run everything as Trigger.dev tasks",
    "Full observability in the dashboard",
], l=Inches(7.2), t=Inches(3.2), w=Inches(5), size=17)

# Decision callout
rect(s, Inches(1.2), Inches(5.3), Inches(11), Inches(1.2), GRAY_3)
text(s, Inches(1.5), Inches(5.45), Inches(10.5), Inches(0.9),
     "Decision: Option B \u2014 demonstrates Trigger.dev's strengths:\n"
     "long-running tasks, observability, background processing, and chat.agent() streaming",
     size=18, color=GREEN_MID)


# ── 5. DRY RUN ───────────────────────────────────────────────────
s = blank()
section_header(s, "The Dry Run", "First contact with the prerelease SDK", 5)

bullets(s, [
    "1.  npx create-next-app@latest \u2014 Next.js 16, React 19, Turbopack",
    "2.  npm install @trigger.dev/sdk@0.0.0-chat-prerelease-20260415164455",
    "3.  npx trigger dev \u2014 initialized project, got secret keys",
    "4.  Installed AI Elements components (conversation, message, shimmer ...)",
    "5.  Installed LanceDB, AI SDK, OpenRouter adapter",
    "",
    "First surprises:",
    '   CLI binary is "trigger", not "trigger.dev" (docs mismatch)',
    "   task() doesn't accept schema \u2014 need schemaTask()",
    "   Next.js 16 APIs differ from training data",
    "   LanceDB native binary can't be bundled by Turbopack",
], t=Inches(2.5))


# ── 6. ARCHITECTURE ──────────────────────────────────────────────
s = blank()
section_header(s, "Architecture", "End-to-end system design", 6)

arch = (
    "User Question\n"
    "     \u2502\n"
    "     \u25bc\n"
    "Next.js Frontend  (useTriggerChatTransport)\n"
    "     \u2502\n"
    "     \u25bc  getChatToken() server action\n"
    "     \u2502\n"
    "     \u25bc\n"
    "Trigger.dev Worker:  chat.agent('trigger-support')\n"
    "     \u2502\n"
    "     \u251c\u2500\u2500\u2500\u25b6  expandQuery()   \u2014 Gemini rewrites to SDK terms\n"
    "     \u251c\u2500\u2500\u2500\u25b6  runSearch()     \u2014 Hybrid: Vector + BM25 + RRF\n"
    "     \u251c\u2500\u2500\u2500\u25b6  buildContext()  \u2014 Format chunks as Markdown\n"
    "     \u251c\u2500\u2500\u2500\u25b6  streamText()   \u2014 Gemini 3 Flash streams answer\n"
    "     \u2502\n"
    "     \u25bc\n"
    "Stream to UI \u2014 Citation badges + Sources list"
)
text(s, Inches(2), Inches(2.3), Inches(9), Inches(5),
     arch, size=16, color=GREEN_1, font=FONT_MONO)


# ── 7. TECH STACK ────────────────────────────────────────────────
s = blank()
section_header(s, "Tech Stack", n=7)

stack = [
    ("Framework",   "Next.js 16 (App Router, React 19, Turbopack)"),
    ("Chat Agent",  "@trigger.dev/sdk chat.agent() prerelease"),
    ("LLM",         "Google Gemini 3 Flash via OpenRouter"),
    ("Embeddings",  "OpenAI text-embedding-3-small (1536 dims)"),
    ("Vector DB",   "LanceDB (embedded Rust binary, worker-only)"),
    ("Search",      "Hybrid BM25 + Vector with RRF Reranking"),
    ("UI",          "AI Elements + shadcn/ui + Streamdown"),
    ("Styling",     "Tailwind CSS v4"),
]

y = Inches(2.4)
for label, value in stack:
    text(s, Inches(1.2), y, Inches(2.5), Inches(0.4),
         label, size=17, color=GRAY_1, bold=True)
    text(s, Inches(3.8), y, Inches(8), Inches(0.4),
         value, size=17, color=WHITE)
    y += Inches(0.55)


# ── 8. PHASES 1-5 ────────────────────────────────────────────────
s = blank()
section_header(s, "Implementation Phases 1\u20135", "Foundation", 8)

bullets(s, [
    "Phase 1: Design Discussion",
    "   Chose RAG over MCP/inline fetch. LanceDB for embedded hybrid search.",
    "",
    "Phase 2: Indexing Pipeline  (trigger/indexDocs.ts)",
    "   GitHub Tree API \u2192 313 MDX files \u2192 section chunking \u2192 embeddings \u2192 LanceDB",
    "",
    "Phase 3: Search Pipeline  (trigger/searchDocs.ts)",
    "   Hybrid vector + BM25 + RRF reranking. schemaTask for dashboard testing.",
    "",
    "Phase 4: Chat Agent  (trigger/chat.ts)",
    "   chat.agent() with streamText. Hit two critical SDK bugs (next slide).",
    "",
    "Phase 5: First Successful Chat",
    "   Pre-fetch pattern working. Gemini streaming answers with inline citations.",
], t=Inches(2.5), size=18)


# ── 9. SDK BUGS ──────────────────────────────────────────────────
s = blank()
section_header(s, "Prerelease SDK Bugs", "Two blockers discovered, both worked around", 9)

# Bug 1
rect(s, Inches(1.2), Inches(2.4), Inches(5.3), Inches(4.0), GRAY_3)
text(s, Inches(1.5), Inches(2.55), Inches(4.8), Inches(0.4),
     "Bug 1: maxSteps broken", size=20, color=RED_SOFT, bold=True)
bullets(s, [
    "chat.agent() fires turn-complete",
    "after first LLM finish event",
    "",
    "Stream closes with finishReason:",
    '"tool-calls" \u2014 no text delivered',
    "",
    "Fix: Pre-fetch search results,",
    "inject into system prompt,",
    "single-step stream",
], l=Inches(1.5), t=Inches(3.1), w=Inches(4.8), size=16)

# Bug 2
rect(s, Inches(6.9), Inches(2.4), Inches(5.3), Inches(4.0), GRAY_3)
text(s, Inches(7.2), Inches(2.55), Inches(4.8), Inches(0.4),
     "Bug 2: ai.toolExecute() hangs", size=20, color=RED_SOFT, bold=True)
bullets(s, [
    "Parent run suspended at",
    "Waitpoint waiting for webhook",
    "that never fires",
    "",
    "Run stuck in 'Suspended' state",
    "indefinitely",
    "",
    "Fix: Call runSearch() directly",
    "as a plain async function",
], l=Inches(7.2), t=Inches(3.1), w=Inches(4.8), size=16)


# ── 10. PHASES 6-10 ──────────────────────────────────────────────
s = blank()
section_header(s, "Implementation Phases 6\u201310", "UI, polish & documentation", 10)

bullets(s, [
    "Phase 6: Citation UI",
    "   Inline badges with hover cards, regex parsing of ([label](url)) pattern",
    "",
    "Phase 7: UI Polish",
    "   Sources panel, layout scroll fix, header, Triggerito mascot, Streamdown fix",
    "",
    "Phase 8: Query Expansion",
    "   expandQuery() \u2014 Gemini rewrites natural language to SDK terms",
    "   10 domain-specific term mappings for common support topics",
    "",
    "Phase 9: Golden QA Evaluation Framework",
    "   6-question test suite with verified golden answers and scoring rubric",
    "",
    "Phase 10: Documentation",
    "   README, GUIDE (step-by-step build), AGENTS.md, evaluation reports",
], t=Inches(2.5), size=18)


# ── 11. RAG DEEP DIVE ────────────────────────────────────────────
s = blank()
section_header(s, "RAG Pipeline Deep Dive", "How search actually works", 11)

bullets(s, [
    "1.  Query Expansion  (expandQuery)",
    '    "pause a job until human approves"  \u2192  "wait.forToken waitpoint checkpoint"',
    "",
    "2.  Multi-Query Retrieval",
    "    Primary: expanded query (12 results) + Secondary: original query (6 results)",
    "    Merged by chunk ID, primary takes priority, capped at 12",
    "",
    "3.  Hybrid Search  (per query)",
    "    Vector: cosine similarity via text-embedding-3-small",
    "    BM25: full-text on content + section columns",
    "    RRF: Reciprocal Rank Fusion merges both rankings",
    "",
    "4.  Context Injection",
    "    Top chunks formatted as Markdown sections with source URLs",
    "    Prepended to system prompt before streamText()",
], t=Inches(2.5), size=18)


# ── 12. EVAL BASELINE ────────────────────────────────────────────
s = blank()
section_header(s, "Evaluation: Baseline", "First run \u2014 5.8 / 10 overall", 12)

# Results table
rows = [
    ("Q1", "Self-hosted limitations",    "3/10",  "FAIL",  RED_SOFT),
    ("Q2", "triggerAndWait compute time", "9/10",  "PASS",  GREEN_1),
    ("Q3", "Cron every 30 minutes",       "9/10",  "PASS",  GREEN_1),
    ("Q4", "Cost of runs",                "3/10",  "FAIL",  RED_SOFT),
    ("Q5", "Exit task on criteria",       "4/10",  "FAIL",  RED_SOFT),
    ("Q6", "Production secret key",       "7/10",  "PASS",  GREEN_1),
]
y = Inches(2.5)
for q, desc, score, verdict, vcolor in rows:
    text(s, Inches(1.2), y, Inches(0.6), Inches(0.35), q, size=16, color=GRAY_1, bold=True)
    text(s, Inches(1.9), y, Inches(4.5), Inches(0.35), desc, size=16, color=WHITE)
    text(s, Inches(6.5), y, Inches(1), Inches(0.35), score, size=16, color=vcolor, bold=True)
    text(s, Inches(7.6), y, Inches(1.5), Inches(0.35), verdict, size=16, color=vcolor)
    y += Inches(0.45)

bullets(s, [
    "Root causes:",
    "   Chunks too large (2500 chars) \u2014 tables and short sections dropped",
    "   Pre-heading content silently lost",
    "   Query expansion had no domain vocabulary",
    "   Retrieval depth too shallow (top 8)",
], l=Inches(1.2), t=Inches(5.4), w=Inches(10), size=16)


# ── 13. ROUND 2 FIXES ────────────────────────────────────────────
s = blank()
section_header(s, "Round 2: Fixes Applied", "Chunking, expansion, retrieval depth", 13)

bullets(s, [
    "1.  Improved Chunking  (indexDocs.ts)",
    "    Cap: 2500 \u2192 1500 chars",
    "    Added 200-char overlapping sub-chunks",
    "    Min body: 80 \u2192 40 chars (captures tables and short sections)",
    "    Pre-heading intro content now captured as 'Overview' chunk",
    "",
    "2.  Query Expansion Vocabulary  (chat.ts)",
    "    10 domain mappings: cost \u2192 TRQL/compute_cost, exit \u2192 AbortTaskRunError, etc.",
    "",
    "3.  Retrieval Depth  (searchDocs.ts)",
    "    limit: 8 \u2192 12 results",
    "",
    "Re-indexed: 313 files \u2192 1,952 chunks  (up from ~800)",
], t=Inches(2.5), size=18)


# ── 14. ROUND 2 RESULTS ──────────────────────────────────────────
s = blank()
section_header(s, "Round 2: Results", "5.8 \u2192 8.8 / 10  (+3.0)", 14)

rows2 = [
    ("Q1", "Self-hosted limitations",    "3", "8", "+5", "Warm starts, auto-scaling, checkpoints surfaced"),
    ("Q2", "triggerAndWait compute",      "9", "9", "+0", "Now uses 'checkpointed' term"),
    ("Q3", "Cron every 30 min",           "9", "9", "+0", "Both declarative and imperative approaches"),
    ("Q4", "Cost of runs",                "3", "9", "+6", "Full TRQL SQL query with compute_cost"),
    ("Q5", "Exit task on criteria",       "4", "9", "+5", "Leads with AbortTaskRunError + code"),
    ("Q6", "Secret key",                  "7", "9", "+2", "Includes tr_prod_ prefix"),
]
y = Inches(2.5)
text(s, Inches(1.2), y, Inches(0.6), Inches(0.35), "", size=14, color=GRAY_2)
text(s, Inches(1.9), y, Inches(3.5), Inches(0.35), "", size=14, color=GRAY_2)
text(s, Inches(5.5), y, Inches(0.8), Inches(0.35), "Before", size=13, color=GRAY_2, bold=True)
text(s, Inches(6.4), y, Inches(0.8), Inches(0.35), "After", size=13, color=GRAY_2, bold=True)
text(s, Inches(7.3), y, Inches(0.6), Inches(0.35), "\u0394", size=13, color=GRAY_2, bold=True)
text(s, Inches(8.0), y, Inches(5), Inches(0.35), "Key improvement", size=13, color=GRAY_2, bold=True)
y += Inches(0.4)

for q, desc, before, after, delta, note in rows2:
    dcolor = GREEN_1 if delta != "+0" else GRAY_2
    text(s, Inches(1.2), y, Inches(0.6), Inches(0.35), q, size=15, color=GRAY_1, bold=True)
    text(s, Inches(1.9), y, Inches(3.5), Inches(0.35), desc, size=15, color=WHITE)
    text(s, Inches(5.7), y, Inches(0.5), Inches(0.35), before, size=15, color=GRAY_1)
    text(s, Inches(6.5), y, Inches(0.5), Inches(0.35), after, size=15, color=GREEN_1, bold=True)
    text(s, Inches(7.3), y, Inches(0.6), Inches(0.35), delta, size=15, color=dcolor, bold=True)
    text(s, Inches(8.0), y, Inches(5), Inches(0.35), note, size=14, color=GRAY_1)
    y += Inches(0.45)


# ── 15. ROUND 3 FIXES ────────────────────────────────────────────
s = blank()
section_header(s, "Round 3: Fine-Tuning", "Five additional improvements", 15)

bullets(s, [
    "1.  Page Context Prefix",
    '    Each chunk: "Page: self-hosting/overview | Section: Feature Comparison"',
    "",
    "2.  Semantic JSX Tag Preservation",
    "    <Note>, <Tip>, <Warning> \u2192 **Note:** instead of stripped",
    "",
    "3.  Multi-Query Retrieval",
    "    Runs BOTH expanded + original query, merges + deduplicates",
    "",
    "4.  BM25 Index on 'section' Column",
    '    Heading-level matching: "declarative schedules" hits section titles',
    "",
    "5.  System Prompt Tuning",
    "    Always include env var names, key prefixes, numeric limits",
], t=Inches(2.5), size=18)


# ── 16. FINAL RESULTS ────────────────────────────────────────────
s = blank()
section_header(s, "Final Results: 9.3 / 10", "Three rounds of improvement", 16)

rows3 = [
    ("Q1", "Self-hosted limits",  "3",  "8",  "9",  GREEN_1),
    ("Q2", "Wait compute time",   "9",  "9",  "9",  WHITE),
    ("Q3", "Cron every 30 min",    "9",  "9", "10",  GREEN_1),
    ("Q4", "Cost of runs",         "3",  "9",  "9",  GREEN_1),
    ("Q5", "Exit task",            "4",  "9", "10",  GREEN_1),
    ("Q6", "Secret key",           "7",  "9",  "9",  GREEN_1),
]

# Header row
y = Inches(2.4)
for col, lbl in [(Inches(5.5), "Baseline"), (Inches(6.8), "Round 2"), (Inches(8.1), "Round 3")]:
    text(s, col, y, Inches(1.2), Inches(0.35), lbl, size=14, color=GRAY_2, bold=True)
y += Inches(0.4)
# Divider
rect(s, Inches(1.2), y, Inches(8.5), Pt(1), GRAY_3)
y += Inches(0.15)

for q, desc, b, r2, r3, r3c in rows3:
    text(s, Inches(1.2), y, Inches(0.6), Inches(0.35), q, size=16, color=GRAY_1, bold=True)
    text(s, Inches(1.9), y, Inches(3.5), Inches(0.35), desc, size=16, color=WHITE)
    text(s, Inches(5.7), y, Inches(0.7), Inches(0.35), b, size=16, color=GRAY_1)
    text(s, Inches(7.0), y, Inches(0.7), Inches(0.35), r2, size=16, color=WHITE)
    text(s, Inches(8.3), y, Inches(0.7), Inches(0.35), r3, size=16, color=r3c, bold=True)
    y += Inches(0.45)

# Overall callout
rect(s, Inches(1.2), y + Inches(0.15), Inches(8.5), Inches(0.7), GRAY_3)
text(s, Inches(1.5), y + Inches(0.22), Inches(8), Inches(0.55),
     "Overall:  5.8  \u2192  8.8  \u2192  9.3 / 10", size=24, color=GREEN_1, bold=True)


# ── 17. PRODUCTION ───────────────────────────────────────────────
s = blank()
section_header(s, "Production Deployment", "Live at triggerdev-support.vercel.app", 17)

bullets(s, [
    "Frontend:  Vercel (Next.js 16, serverless)",
    "Worker:    Trigger.dev Cloud (tasks, search, chat agent)",
    "Vector DB: LanceDB on worker filesystem",
    "",
    "Key insight:",
    "   Vercel's ephemeral filesystem is fine \u2014 all vector DB ops run",
    "   on the Trigger.dev worker, not on Vercel. API route only triggers",
    "   tasks by string ID.",
    "",
    "Deployment:",
    "   1. Deploy worker to Trigger.dev cloud",
    "   2. Set env vars in Vercel + Trigger.dev dashboard",
    "   3. curl -X POST .../api/index-docs  (one-time)",
    "   4. Ready to chat",
], t=Inches(2.5), size=18)


# ── 18. HONEST GAPS ──────────────────────────────────────────────
s = blank()
section_header(s, "Honest Gaps", "What's NOT production-ready today", 18)

# Three columns
col_data = [
    ("Search & Retrieval", RED_SOFT, [
        "Local LanceDB on single worker",
        "No shared state across instances",
        "No auto re-index on doc changes",
        "MDX tables still chunk poorly",
        "6-question eval is too small",
    ]),
    ("Agent Intelligence", AMBER, [
        "No multi-turn conversation",
        "No tool calling (can't look up",
        "  a user's actual runs/billing)",
        'No confident "I don\'t know"',
        "No conversation memory",
    ]),
    ("Operational", PURPLE, [
        "No feedback loop from users",
        "No analytics on question types",
        "No auth or rate limiting",
        "No escalation to humans",
        "No A/B testing of prompts",
    ]),
]

for i, (title, title_color, items) in enumerate(col_data):
    x = Inches(1.2) + Inches(3.8) * i
    rect(s, x, Inches(2.4), Inches(3.5), Inches(4.3), GRAY_3)
    text(s, x + Inches(0.3), Inches(2.55), Inches(3), Inches(0.4),
         title, size=18, color=title_color, bold=True)
    bullets(s, items, l=x + Inches(0.3), t=Inches(3.1), w=Inches(3), size=15,
            color=GRAY_1, line_spacing=Pt(7))


# ── 19. ROADMAP PHASE 1 ──────────────────────────────────────────
s = blank()
section_header(s, "Roadmap: Make It Reliable", "Weeks 1\u20134", 19)

bullets(s, [
    "1.  Multi-step tool calling  (once SDK ships maxSteps fix)",
    "    Agent decides WHEN and WHAT to search, can refine queries",
    "    Observable tool calls in the Trigger.dev dashboard",
    "",
    "2.  Conversation memory & follow-ups",
    "    Track state across turns, resolve references to prior answers",
    "",
    "3.  Confidence scoring & \"I don't know\"",
    "    Low retrieval score \u2192 acknowledge gap, offer human escalation or Discord link",
    "",
    "4.  Automated eval pipeline",
    "    50+ golden QA pairs, run on every PR via Trigger.dev scheduled task",
    "    Regression detection: alert if any score drops",
], t=Inches(2.5), size=18)


# ── 20. ROADMAP PHASE 2 ──────────────────────────────────────────
s = blank()
section_header(s, "Roadmap: Make It Smart", "Months 2\u20133", 20)

bullets(s, [
    "5.  Live data tools  (not just docs)",
    "    Connect to Trigger.dev API: look up user's runs, errors, billing",
    "    Connect to GitHub: search issues, PRs, changelogs",
    '    "Why is my task failing?" \u2192 agent checks their actual run logs',
    "",
    "6.  Auto-reindex on docs changes",
    "    GitHub webhook triggers incremental re-indexing task",
    "",
    "7.  Source code search",
    "    Index the open-source repo alongside docs",
    '    "How does queue concurrency work internally?" \u2192 actual source code',
    "",
    "8.  Feedback loop",
    "    Thumbs up/down on every answer, stored in DB",
    "    Weekly report: worst questions \u2192 targeted RAG improvements",
], t=Inches(2.5), size=18)


# ── 21. ROADMAP PHASE 3 ──────────────────────────────────────────
s = blank()
section_header(s, "Roadmap: Make It a Product", "Months 3\u20136", 21)

bullets(s, [
    "9.   Human-in-the-loop escalation",
    "     Agent recognizes it can't help \u2192 creates support ticket",
    "     Passes full conversation context (no repeat explaining)",
    "",
    "10.  Multi-channel deployment",
    "     Dashboard sidebar widget, Discord bot, GitHub Discussions auto-responder",
    "",
    "11.  Personalization & auth",
    "     Auth-aware: knows user's project, plan tier, recent errors",
    "     Proactive: \"I see your task crashed 3x today \u2014 here's why\"",
    "",
    "12.  Production infrastructure",
    "     Remote vector DB (Pinecone/Qdrant) for multi-worker scaling",
    "     Rate limiting, abuse detection, cost controls",
    "     A/B test models, prompts, retrieval strategies",
], t=Inches(2.5), size=18)


# ── 22. DX FEEDBACK ──────────────────────────────────────────────
s = blank()
section_header(s, "Developer Experience Feedback", "Trigger.dev platform & AI Agents SDK", 22)

# Positives
text(s, Inches(1.2), Inches(2.4), Inches(5.3), Inches(0.4),
     "What worked well", size=20, color=GREEN_1, bold=True)
bullets(s, [
    "Dashboard is excellent \u2014 stream viewer,",
    "  run inspector, logs are first-class",
    "chat.agent() API is clean for the happy path",
    "useTriggerChatTransport: seamless frontend",
    "Task observability is a real differentiator",
    "schemaTask for typed inputs is great",
    "Re-running tasks from dashboard: very useful",
], l=Inches(1.2), t=Inches(3.0), w=Inches(5), size=16, color=WHITE)

# Negatives
text(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(0.4),
     "What needs work", size=20, color=RED_SOFT, bold=True)
bullets(s, [
    "maxSteps broken in chat.agent()",
    "  #1 blocker for real agent workflows",
    "ai.toolExecute() hangs indefinitely",
    "TypeScript types lag behind runtime",
    "  (ignoreBuildErrors required)",
    "Prerelease version pinning is fragile",
    "CLI binary name mismatch vs. docs",
    "No docs mention the maxSteps limitation",
], l=Inches(7.2), t=Inches(3.0), w=Inches(5), size=16, color=WHITE)


# ── 23. KEY TAKEAWAY ─────────────────────────────────────────────
s = blank()
set_bg(s)
logo(s, Inches(11.5), Inches(0.35), Inches(0.35))
slide_num(s, 23)

text(s, Inches(1.2), Inches(1.0), Inches(10), Inches(0.8),
     "Key Takeaway", size=38, color=WHITE, bold=True)
rect(s, Inches(1.2), Inches(1.85), Inches(1.8), Pt(3), GREEN_1)

text(s, Inches(1.2), Inches(2.3), Inches(10), Inches(1.0),
     "The MVP proves the pattern works:", size=22, color=WHITE)

# Highlighted formula
rect(s, Inches(1.2), Inches(3.1), Inches(10.5), Inches(0.8), GRAY_3)
text(s, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
     "chat.agent()  +  hybrid RAG  +  streaming UI  =  accurate, cited answers in < 3 seconds",
     size=20, color=GREEN_1, bold=True, font=FONT_MONO)

text(s, Inches(1.2), Inches(4.3), Inches(10), Inches(0.8),
     "The gap to world-class is not better search or a smarter model.",
     size=20, color=WHITE)

text(s, Inches(1.2), Inches(5.0), Inches(10), Inches(0.5),
     "It's the feedback loop:", size=22, color=GREEN_MID, bold=True)

bullets(s, [
    "Know which questions users actually ask",
    "Know which answers are wrong",
    "Know when to escalate to a human",
    "Learn from every interaction",
], l=Inches(1.5), t=Inches(5.5), w=Inches(10), size=18, color=WHITE, line_spacing=Pt(6))


# ── 24. THANK YOU ────────────────────────────────────────────────
s = blank()
set_bg(s)

# Top accent
rect(s, 0, 0, Inches(8), Pt(4), GREEN_1)
rect(s, Inches(8), 0, Inches(5.333), Pt(4), GREEN_2)

# Logo
s.shapes.add_picture("docs/logo.png", Inches(1.2), Inches(1.6), height=Inches(0.6))

text(s, Inches(1.2), Inches(2.7), Inches(10), Inches(1.0),
     "Thank You", size=48, color=WHITE, bold=True)

# URL in green
text(s, Inches(1.2), Inches(3.8), Inches(10), Inches(0.6),
     "triggerdev-support.vercel.app", size=24, color=GREEN_1)

# Contact
text(s, Inches(1.2), Inches(4.8), Inches(10), Inches(1.5),
     "Juan Suarez\njuanluis.suarez@gmail.com", size=18, color=GRAY_1)

# Mascot
s.shapes.add_picture("docs/triggerito.png", Inches(10), Inches(2.5), height=Inches(2.8))

# Bottom accent
rect(s, 0, H - Pt(4), Inches(5), Pt(4), BLUE)
rect(s, Inches(5), H - Pt(4), Inches(8.333), Pt(4), PURPLE)

slide_num(s, 24)


# ── SAVE ──────────────────────────────────────────────────────────
output = "docs/Trigger-Support-Presentation.pptx"
prs.save(output)
print(f"Saved to {output}")
